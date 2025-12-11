import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

RAW_DATA_PATH = "data/raw"
TEXT_OUTPUT_PATH = "data/text"

os.makedirs(TEXT_OUTPUT_PATH, exist_ok=True)

def extract_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_from_pptx(file_path):
    ppt = Presentation(file_path)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_documents():
    for filename in os.listdir(RAW_DATA_PATH):
        file_path = os.path.join(RAW_DATA_PATH, filename)

        if filename.endswith(".pdf"):
            text = extract_from_pdf(file_path)

        elif filename.endswith(".docx"):
            text = extract_from_docx(file_path)

        elif filename.endswith(".pptx"):
            text = extract_from_pptx(file_path)

        elif filename.endswith(".txt"):
            text = extract_from_txt(file_path)

        else:
            print(f"Skipped unsupported file: {filename}")
            continue

        output_file = os.path.join(TEXT_OUTPUT_PATH, filename + ".txt")

        with open(output_file, "w", encoding="utf-8") as f:
            f.write(text)

        print(f"Extracted: {filename}")

if __name__ == "__main__":
    extract_documents()
